import streamlit as st
from supabase import create_client, Client
from dotenv import load_dotenv
from urllib.parse import quote
from datetime import datetime
import os
import re
import unicodedata

#processing imports
from io import BytesIO
import fitz  # PyMuPDF
import pandas as pd
import pytesseract
from PIL import Image
from pptx import Presentation


pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'


def sanitize_filename(name):
    name = unicodedata.normalize('NFKD', name).encode('ascii', 'ignore').decode()
    name = name.replace(' ', '_')
    name = re.sub(r'[^\w\-_\.]', '', name)
    return name

def extract_text_from_pdf(file_bytes):
    doc = fitz.open(stream=BytesIO(file_bytes), filetype="pdf")
    return "\n".join([page.get_text() for page in doc])

def extract_text_from_excel(file_bytes):
    df = pd.read_excel(BytesIO(file_bytes))
    return df.to_string(index=False)

def extract_text_from_csv(file_bytes):
    df = pd.read_csv(BytesIO(file_bytes))
    return df.to_string(index=False)

def extract_text_from_image(file_bytes):
    try:
        image = Image.open(BytesIO(file_bytes))
        text = pytesseract.image_to_string(image)
        return text
    except Exception as e:
        return f"Image text extraction failed: {str(e)}"

def extract_text_from_ppt(file_bytes):
    try:
        prs = Presentation(BytesIO(file_bytes))
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"PPT text extraction failed: {str(e)}"


load_dotenv()

SUPABASE_URL = os.getenv("SUPABASE_URL")
SUPABASE_KEY = os.getenv("SUPABASE_KEY")
BUCKET_NAME = os.getenv("SUPABASE_BUCKET")

supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

#UI setup

st.title("📤 Upload File & Extract Text")

uploaded_file = st.file_uploader("Choose a file", type=["pdf", "csv", "xlsx", "pptx", "png", "jpg", "jpeg"])

if uploaded_file is not None:
    file_bytes = uploaded_file.read()
    original_filename = uploaded_file.name
    safe_filename = sanitize_filename(original_filename)
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    supabase_path = f"{timestamp}_{safe_filename}"

    extracted_text = ""
    filetype = uploaded_file.type

    try:
        if filetype == "application/pdf":
            extracted_text = extract_text_from_pdf(file_bytes)
        elif filetype == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
            extracted_text = extract_text_from_excel(file_bytes)
        elif filetype == "text/csv":
            extracted_text = extract_text_from_csv(file_bytes)
        elif filetype == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            extracted_text = extract_text_from_ppt(file_bytes)
        elif filetype in ["image/jpeg", "image/png"]:
            extracted_text = extract_text_from_image(file_bytes)
        else:
            extracted_text = "Unsupported file type for text extraction."
    except Exception as e:
        extracted_text = f"Text extraction failed: {str(e)}"

    try:
        response = supabase.storage.from_(BUCKET_NAME).upload(
            supabase_path,
            file_bytes,
            {"content-type": uploaded_file.type}
        )
        if response:
            raw_url = supabase.storage.from_(BUCKET_NAME).get_public_url(supabase_path)
            public_url = raw_url.rsplit("/", 1)[0] + "/" + quote(supabase_path)
            st.success(" File uploaded to Supabase Storage!")
            st.markdown(f"[🔗 View File]({public_url})")
        else:
            st.error("❌ Upload failed.")
    except Exception as e:
        st.error(f"❌ Error during file upload: {str(e)}")

    try:
        db_response = supabase.table("extracted_texts").insert({
            "filename": safe_filename,
            "filetype": filetype,
            "content": extracted_text
        }).execute()
        st.success("📝 Extracted text saved to Supabase database!")
    except Exception as e:
        st.error(f"⚠️ Failed to save extracted text: {str(e)}")

    #preview extracted text
    with st.expander("📄 Preview Extracted Text"):
        st.text_area("Extracted Text", extracted_text, height=300)
